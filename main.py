import os
import shutil
from pptx import Presentation
from pptx.util import Inches
from pytube import YouTube
import cv2
from skimage.metrics import structural_similarity

#############################################################
default_tmp_folder = "tmp"
default_frames_per_second_to_extract = 0.5
default_time_period=[(0, 36000)]

# Get user input
url = input("Enter the YouTube video URL: ")
tmp_folder = input(f"Enter the temporary folder path (default: {default_tmp_folder}): ") or default_tmp_folder
frames_per_second_to_extract = input(f"Enter the frames per second to extract (default: {default_frames_per_second_to_extract}): ") or default_frames_per_second_to_extract

time_periods = []
time_period_str = input(f"Enter time periods (default: {default_time_period}) or leave empty to use default: ")
if time_period_str:
    try:
        time_periods = eval(time_period_str)
    except Exception as e:
        print("Error parsing time periods:", e)
        time_periods = default_time_period
else:
    time_periods = default_time_period

debug = input(f"Enable debug mode? (y/n, default: No): ").lower() == "y"
#############################################################


def down_yt(url, output_path, filename):
    try:
        yt = YouTube(url)
        stream = yt.streams.get_highest_resolution()
        print(f"Downloading '{yt.title}'...")
        stream.download(output_path, filename=filename)
        print("Download completed!")
    except Exception as e:
        print(f"An error occurred: {e}")

def compile_to_ppt(image_folder, output_ppt):
    # Create a PowerPoint presentation
    prs = Presentation()
    image_file = sorted(os.listdir(image_folder))
    image_file = image_file[1:]
    frame_numbers = [int(file.split("_")[1].split(".")[0]) for file in image_file]
    image_files = [f"frame_{frame_num}.jpg" for frame_num in sorted(frame_numbers)]
    # Iterate through each image file in the folder
    for image_file in image_files:
        if image_file.endswith(".jpg") or image_file.endswith(".jpeg"):
            image_path = os.path.join(image_folder, image_file)

            # Add a slide to the presentation
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Using blank slide layout

            # Add the image to the slide
            left = top = Inches(0)
            slide.shapes.add_picture(image_path, left, top, width=prs.slide_width, height=prs.slide_height)

    # Save the PowerPoint presentation
    prs.save(output_ppt)
    print(f"Presentation saved as '{output_ppt}'.")

def extract_frames(video_path, output_folder, fps, time_periods=[(0, 3600)], mode="whitelist"):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    video_capture = cv2.VideoCapture(video_path)
    frame_rate = int(video_capture.get(cv2.CAP_PROP_FPS))
    total_frames = int(video_capture.get(cv2.CAP_PROP_FRAME_COUNT))
    frame_interval = int(frame_rate / fps)
    print(f"frame rate ={frame_rate}")
    if mode == 'whitelist':
        for start, end in time_periods:
            # video_capture.set(cv2.CAP_PROP_POS_MSEC, start * 1000)
            video_capture.set(cv2.CAP_PROP_POS_FRAMES, start * frame_rate)
            frame_count = int(video_capture.get(cv2.CAP_PROP_POS_FRAMES))
            print(f"setting frame count to {frame_count}")
            if end > total_frames / frame_rate:
                end = total_frames / frame_rate

            while True:
                # Read the next frame
                success, frame = video_capture.read()

                if not success:
                    print(f"read error at {frame_count}")
                    break

                # Calculate current time in seconds
                current_time_seconds = frame_count / frame_rate
                if current_time_seconds > end:
                    print("out of time")
                    break
                if frame_count % frame_interval == 0:
                    # Save frame
                    frame_filename = os.path.join(output_folder, f"frame_{frame_count}.jpg")
                    print(f"saving {frame_count}")
                    cv2.imwrite(frame_filename, frame)

                frame_count += 1
                print(f"currently at {frame_count}")

    video_capture.release()


def is_different(frame1_path, frame2_path, threshold=0.4):
    # Read images
    frame1 = cv2.imread(frame1_path)
    frame2 = cv2.imread(frame2_path)

    # Convert images to grayscale
    gray_frame1 = cv2.cvtColor(frame1, cv2.COLOR_BGR2GRAY)
    gray_frame2 = cv2.cvtColor(frame2, cv2.COLOR_BGR2GRAY)

    # Calculate Structural Similarity Index (SSI)
    (score, diff) = structural_similarity(gray_frame1, gray_frame2, full=True)
    #print(f"file1 {frame1_path}, file2 {frame2_path}, score{score}")

    # Calculate Mean Squared Error (MSE)
    mse = ((gray_frame1 - gray_frame2) ** 2).mean()

    # Check if the frames are more than threshold different
    return score < threshold


def find_largest_frame_number(folder_path):
    # Get list of frame files
    frame_files = os.listdir(folder_path)

    # Extract frame numbers from filenames
    frame_numbers = [int(filename.split('_')[1].split('.')[0]) for filename in frame_files]

    # Find the largest frame number
    largest_frame_number = max(frame_numbers)

    return largest_frame_number


def mv_diff_frames(frame_folder, threshold=0.8):
    # Create a folder for changed frames
    changed_frames_folder = os.path.join(frame_folder, "changed_frames")
    if not os.path.exists(changed_frames_folder):
        os.makedirs(changed_frames_folder)

    # Get list of frame files
    frame_files = sorted(os.listdir(frame_folder))
    frame_files = frame_files[1:]
    frame_numbers = [int(file.split("_")[1].split(".")[0]) for file in frame_files]
    frame_files = [f"frame_{frame_num}.jpg" for frame_num in sorted(frame_numbers)]

    print("framefiles",frame_files)

    # Compare consecutive frames
    for i in range(len(frame_files) - 1):
        frame1_path = os.path.join(frame_folder, frame_files[i])
        frame2_path = os.path.join(frame_folder, frame_files[i + 1])
        #print(frame1_path,frame2_path)

        if is_different(frame1_path, frame2_path, threshold):
            # Copy the first frame to the changed frames folder
            changed_frame_path = os.path.join(changed_frames_folder, frame_files[i])
            shutil.copyfile(frame1_path, changed_frame_path)
    last_frame_path = os.path.join(frame_folder, frame_files[-1])
    changed_last_frame_path = os.path.join(changed_frames_folder, frame_files[-1])
    shutil.copyfile(last_frame_path, changed_last_frame_path)



try:
    shutil.rmtree(tmp_folder)
    print(f"Directory '{tmp_folder}' removed successfully.")
except OSError as e:
    print(f"Error: {tmp_folder} : {e.strerror}")

down_yt(url,tmp_folder,filename="tmp.mp4")
print(f"whitelisted time {time_periods}")
extract_frames("tmp/tmp.mp4", "tmp/pic", fps=frames_per_second_to_extract, time_periods=time_periods)
#max_frame = find_largest_frame_number("tmp/pic")
#print(max_frame)
mv_diff_frames("tmp/pic")
compile_to_ppt("tmp/pic/changed_frames", "out.ppt")
if not debug:
    try:
        shutil.rmtree(tmp_folder)
        print(f"Directory '{tmp_folder}' removed successfully.")
    except OSError as e:
        print(f"Error: {tmp_folder} : {e.strerror}")